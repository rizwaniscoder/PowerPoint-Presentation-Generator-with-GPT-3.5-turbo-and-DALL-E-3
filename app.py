import streamlit as st
import base64
import pptx
from pptx.util import Inches, Pt
import os
import requests
from io import BytesIO
from PIL import Image
from openai import OpenAI
from dotenv import load_dotenv
load_dotenv()


client = OpenAI()
# OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
client.api_key = os.getenv('OPENAI_API_KEY')


# Constants

GPT_MODEL = "gpt-3.5-turbo"
DALLE_MODEL = "dall-e-3"
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(topic, num_slides):
    system_message = {"role": "system", "content": "You are a helpful assistant."}
    user_message = {"role": "user", "content": f"Generate {num_slides} slide titles for the topic '{topic}'."}

    response = client.chat.completions.create(
        model=GPT_MODEL,
        messages=[system_message, user_message],
    )

    return response.choices[0].message.content.split("\n")[:num_slides]

def generate_image(prompt):
    response = client.images.generate(
        model="dall-e-3",
        prompt=prompt,
        size="1024x1024",
        quality="standard",
        n=1,
        )

    image_url = response.data[0].url
    return image_url

def generate_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'."
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt},
        ],
    )
    return response.choices[0].message.content


def create_presentation(topic, slide_titles, slide_contents, num_slides):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for i in range(num_slides):
        slide_title = slide_titles[i]
        slide_content = slide_contents[i]
        
        # Generate image for the slide
        image_prompt = f"An image related to the slide title: '{slide_title}'."
        image_url = generate_image(image_prompt)

        # Download the image from the URL
        response = requests.get(image_url)
        image = Image.open(BytesIO(response.content))

        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Add image to the slide
        left = Inches(1)  # Adjust the position as needed
        top = Inches(2)
        height = Inches(3)
        pic = slide.shapes.add_picture(BytesIO(response.content), left, top, height=height)

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    prs.save(f"generated_ppt/{topic}_presentation.pptx")


def main():
    st.title("PowerPoint Presentation Generator with GPT-3.5-turbo and DALL·E 3")

    topic = st.text_input("Enter the topic for your presentation:")
    num_slides = st.number_input("Enter the number of slides for the presentation:", min_value=1, value=5, step=1)
    generate_button = st.button("Generate Presentation")

    if generate_button and topic:
        st.info("Generating presentation... Please wait.")
        try:
            slide_titles = generate_slide_titles(topic, num_slides)
            filtered_slide_titles = [item.strip() for item in slide_titles if item.strip()]
            slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
            create_presentation(topic, filtered_slide_titles, slide_contents, num_slides)
            st.success("Presentation generated successfully!")
            st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")

def get_ppt_download_link(topic):
    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"

    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'

if __name__ == "__main__":
    main()
